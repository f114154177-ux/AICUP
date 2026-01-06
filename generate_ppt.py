from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    return slide

def create_content_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    return slide

def add_text_box(slide, left, top, width, height, text, font_size=18):
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = text_box.text_frame
    text_frame.text = text
    text_frame.paragraphs[0].font.size = Pt(font_size)
    return text_box

def add_image(slide, image_path, left, top, width=None, height=None):
    if width:
        slide.shapes.add_picture(image_path, Inches(left), Inches(top), width=Inches(width))
    elif height:
        slide.shapes.add_picture(image_path, Inches(left), Inches(top), height=Inches(height))
    else:
        slide.shapes.add_picture(image_path, Inches(left), Inches(top))

def main():
    # 統計數據
    total_images = 16863
    positive_images = 2787
    negative_images = 14076
    total_bboxes = 2787
    bbox_width_min = 0.01
    bbox_width_avg = 0.0704
    bbox_width_max = 0.18
    bbox_height_min = 0.012
    bbox_height_avg = 0.0852
    bbox_height_max = 0.164
    
    # 創建 PPT
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: 標題頁
    slide1 = create_title_slide(prs, 
                                "AI CUP Training Dataset Analysis",
                                "Statistical Report on Medical Image Dataset")
    
    # Slide 2: Dataset Overview
    slide2 = create_content_slide(prs, "Dataset Overview")
    
    overview_text = f"""
    • Total Images: {total_images:,}
    
    • Training Patients: 50 (patient0001 ~ patient0050)
    
    • Average Images per Patient: ~337
    
    • Image Format: PNG
    
    • Annotation Format: YOLO format (normalized coordinates)
    """
    
    add_text_box(slide2, 0.8, 1.5, 8.5, 4, overview_text, 24)
    
    # Slide 3: Positive vs Negative Data
    slide3 = create_content_slide(prs, "Positive vs Negative Data")
    
    pos_neg_text = f"""
    • Positive Images (with annotations): {positive_images:,}
    
    • Negative Images (without annotations): {negative_images:,}
    
    • Positive Rate: {positive_images/total_images*100:.2f}%
    
    • Total Bounding Boxes: {total_bboxes:,}
    """
    
    add_text_box(slide3, 0.8, 1.5, 4.5, 3, pos_neg_text, 22)
    
    # Add chart if exists
    if os.path.exists("stats_figs/pos_neg_images.png"):
        add_image(slide3, "stats_figs/pos_neg_images.png", 5.5, 1.8, height=4.5)
    
    # Slide 4: Bounding Box Statistics - Width
    slide4 = create_content_slide(prs, "Bounding Box Statistics - Width")
    
    width_stats = f"""
    Width Statistics (Normalized, 0~1):
    
    • Minimum Width: {bbox_width_min:.3f} (~{bbox_width_min*100:.1f}% of image)
    
    • Average Width: {bbox_width_avg:.3f} (~{bbox_width_avg*100:.1f}% of image)
    
    • Maximum Width: {bbox_width_max:.3f} (~{bbox_width_max*100:.1f}% of image)
    """
    
    add_text_box(slide4, 0.8, 1.5, 4.5, 3, width_stats, 20)
    
    if os.path.exists("stats_figs/bbox_width_hist.png"):
        add_image(slide4, "stats_figs/bbox_width_hist.png", 5.5, 1.8, height=4.5)
    
    # Slide 5: Bounding Box Statistics - Height
    slide5 = create_content_slide(prs, "Bounding Box Statistics - Height")
    
    height_stats = f"""
    Height Statistics (Normalized, 0~1):
    
    • Minimum Height: {bbox_height_min:.3f} (~{bbox_height_min*100:.1f}% of image)
    
    • Average Height: {bbox_height_avg:.3f} (~{bbox_height_avg*100:.1f}% of image)
    
    • Maximum Height: {bbox_height_max:.3f} (~{bbox_height_max*100:.1f}% of image)
    """
    
    add_text_box(slide5, 0.8, 1.5, 4.5, 3, height_stats, 20)
    
    if os.path.exists("stats_figs/bbox_height_hist.png"):
        add_image(slide5, "stats_figs/bbox_height_hist.png", 5.5, 1.8, height=4.5)
    
    # Slide 6: Labels per Image Distribution
    slide6 = create_content_slide(prs, "Number of Labels per Image")
    
    labels_text = f"""
    Distribution Analysis:
    
    • Most images have 0 or 1 bounding box
    
    • Total annotated images: {positive_images:,}
    
    • Each positive image has exactly 1 annotation
    
    • Indicates single-target detection task
    """
    
    add_text_box(slide6, 0.8, 1.5, 4.5, 3, labels_text, 20)
    
    if os.path.exists("stats_figs/labels_per_image_hist.png"):
        add_image(slide6, "stats_figs/labels_per_image_hist.png", 5.5, 1.8, height=4.5)
    
    # Slide 7: Spatial Distribution - Where did positive data appear?
    slide7 = create_content_slide(prs, "Where Did Positive Data Appear?")
    
    location_text = f"""
    Spatial Distribution Analysis:
    
    • Heatmap shows bounding box center locations
    
    • X-axis: Horizontal position (0=left, 1=right)
    
    • Y-axis: Vertical position (0=top, 1=bottom)
    
    • Brighter areas indicate more frequent detections
    """
    
    add_text_box(slide7, 0.8, 1.5, 4.5, 3, location_text, 18)
    
    if os.path.exists("stats_figs/positive_center_heatmap.png"):
        add_image(slide7, "stats_figs/positive_center_heatmap.png", 5.5, 1.5, height=5)
    
    # Slide 8: Key Findings Summary
    slide8 = create_content_slide(prs, "Key Findings Summary")
    
    summary_text = """
    1. Dataset Characteristics:
       • Large-scale dataset with 16,863 medical images
       • Highly imbalanced: 16.5% positive, 83.5% negative
    
    2. Bounding Box Properties:
       • Small to medium-sized objects (avg 7% width, 8.5% height)
       • Consistent size distribution across dataset
    
    3. Spatial Patterns:
       • Detections distributed across various locations
       • Heatmap reveals common detection regions
    
    4. Recommendations:
       • Consider data augmentation for imbalanced dataset
       • Use appropriate loss functions for small object detection
       • Apply focal loss or weighted sampling strategies
    """
    
    add_text_box(slide8, 0.8, 1.3, 8.5, 5, summary_text, 16)
    
    # 儲存 PPT
    output_file = "AICUP_Dataset_Analysis_Report.pptx"
    prs.save(output_file)
    print(f"✅ PPT 已成功生成: {output_file}")
    print(f"   共 {len(prs.slides)} 張投影片")

if __name__ == "__main__":
    main()
